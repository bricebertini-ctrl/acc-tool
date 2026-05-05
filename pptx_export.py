"""
Génération du PPTX de proposition commerciale ACC.
Modifie le template_acc.pptx en remplaçant textes, graphiques et tableaux
avec les données calculées par l'application.
"""
from io import BytesIO

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import numpy as np
import pandas as pd

from pptx import Presentation


# ── Helpers texte ──────────────────────────────────────────────────────────────

def _replace_in_shape(shape, replacements: dict):
    """Remplace du texte dans tous les runs d'une shape, en préservant le style."""
    if not hasattr(shape, "text_frame"):
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            for old, new in replacements.items():
                if old in run.text:
                    run.text = run.text.replace(old, str(new))


def _replace_all(slide, replacements: dict):
    for shape in slide.shapes:
        _replace_in_shape(shape, replacements)


def _set_cell_text(cell, text: str):
    """Écrit du texte dans une cellule en conservant le style du premier run."""
    tf = cell.text_frame
    if tf.paragraphs and tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text = text
        for run in tf.paragraphs[0].runs[1:]:
            run.text = ""
    else:
        cell.text = text


# ── Remplacement d'image par position ────────────────────────────────────────

def _swap_picture_at(slide, left_in: float, top_in: float, img_bytes: bytes, tol: float = 0.5):
    """
    Trouve une image proche de (left_in, top_in) et la remplace par img_bytes
    en conservant exactement la même position et la même taille.
    """
    EMU = 914400
    for shape in list(slide.shapes):
        if shape.shape_type != 13:  # PICTURE
            continue
        sl = shape.left  / EMU
        st = shape.top   / EMU
        if abs(sl - left_in) < tol and abs(st - top_in) < tol:
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            shape._element.getparent().remove(shape._element)
            slide.shapes.add_picture(BytesIO(img_bytes), left, top, width, height)
            return True
    return False


# ── Générateurs de graphiques matplotlib ────────────────────────────────────

def _fig_bytes(fig) -> bytes:
    buf = BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight", dpi=150, facecolor="white")
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def chart_prod_mensuelle(monthly: pd.DataFrame, w: float, h: float) -> bytes:
    """Barres mensuelles de production (MWh)."""
    fig, ax = plt.subplots(figsize=(w, h))
    months = monthly.index.tolist()
    vals   = (monthly["Prod (kWh)"] / 1000).tolist()
    ax.bar(months, vals, color="#16a34a", alpha=0.85, edgecolor="white", linewidth=0.5)
    ax.set_ylabel("Production (MWh)", fontsize=10)
    ax.tick_params(axis="x", rotation=45, labelsize=8)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f"{x:,.0f}"))
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.grid(axis="y", alpha=0.3)
    total_gwh = sum(vals) / 1000
    ax.text(0.98, 0.97, f"Total : {total_gwh:.2f} GWh",
            transform=ax.transAxes, ha="right", va="top",
            fontsize=11, fontweight="bold", color="#16a34a")
    fig.tight_layout(pad=0.5)
    return _fig_bytes(fig)


def chart_conso_acc(monthly: pd.DataFrame, w: float, h: float) -> bytes:
    """Barres mensuelles Consommation vs ACC (MWh)."""
    fig, ax = plt.subplots(figsize=(w, h))
    months = monthly.index.tolist()
    conso  = (monthly["Conso (kWh)"] / 1000).tolist()
    acc    = (monthly["ACC (kWh)"]   / 1000).tolist()
    x  = np.arange(len(months))
    bw = 0.38
    ax.bar(x - bw/2, conso, bw, label="Consommation", color="#2563eb", alpha=0.80)
    ax.bar(x + bw/2, acc,   bw, label="ACC",           color="#16a34a", alpha=0.80)
    ax.set_xticks(x)
    ax.set_xticklabels(months, rotation=45, fontsize=8)
    ax.set_ylabel("Énergie (MWh)", fontsize=10)
    ax.legend(fontsize=9, loc="upper right")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.grid(axis="y", alpha=0.3)
    fig.tight_layout(pad=0.5)
    return _fig_bytes(fig)


def chart_taux_acc(monthly: pd.DataFrame, w: float, h: float) -> bytes:
    """Courbes taux d'ACC mensuel (conso et prod)."""
    fig, ax = plt.subplots(figsize=(w, h))
    months     = monthly.index.tolist()
    taux_conso = monthly["Taux ACC conso (%)"].tolist()
    taux_prod  = monthly["Taux ACC prod (%)"].tolist()
    ax.plot(months, taux_conso, "o-",  color="#2563eb", lw=2, ms=5, label="Taux ACC conso")
    ax.plot(months, taux_prod,  "s--", color="#16a34a", lw=2, ms=5, label="Taux ACC prod")
    ax.set_ylabel("Taux d'ACC (%)", fontsize=10)
    ax.tick_params(axis="x", rotation=45, labelsize=8)
    ax.set_ylim(0, 115)
    ax.axhline(100, color="grey", ls=":", alpha=0.5, lw=1)
    ax.legend(fontsize=9)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.grid(axis="y", alpha=0.3)
    fig.tight_layout(pad=0.5)
    return _fig_bytes(fig)


# Couleurs donut (identiques à la charte graphique de la slide)
_DONUT_COMPLEMENT = "#EDE8C4"   # beige clair – Complément réseau
_DONUT_ACC        = "#4472C4"   # bleu – Autoconsommation Collective


def chart_donut_saison(
    acc_kwh: float,
    conso_kwh: float,
    client_name: str,
    saison_label: str,       # "hiver" ou "été"
    w: float,
    h: float,
) -> bytes:
    """
    Donut répartition conso : Complément réseau vs ACC pour une saison.
    Reproduit fidèlement le style de la slide commerciale.
    """
    complement = max(conso_kwh - acc_kwh, 0)
    taux = acc_kwh / conso_kwh * 100 if conso_kwh > 0 else 0

    sizes  = [complement, acc_kwh]
    colors = [_DONUT_COMPLEMENT, _DONUT_ACC]
    labels_raw = ["Complément réseau", "Autoconsommation Collective"]

    fig, ax = plt.subplots(figsize=(w, h))
    fig.patch.set_facecolor("white")

    wedges, _ = ax.pie(
        sizes,
        colors=colors,
        startangle=90,
        counterclock=False,
        wedgeprops=dict(width=0.52, edgecolor="white", linewidth=1.5),
    )

    # Annotations kWh + % à l'extérieur de chaque part
    for i, (wedge, val, lbl) in enumerate(zip(wedges, sizes, labels_raw)):
        angle   = (wedge.theta2 + wedge.theta1) / 2
        rad_out = 0.78
        x = rad_out * np.cos(np.radians(angle))
        y = rad_out * np.sin(np.radians(angle))
        ha = "right" if x < 0 else "left"
        kwh_fmt = f"{val:,.0f} kWh".replace(",", " ")
        pct_fmt = f"{val/conso_kwh*100:.1f} %"
        ax.annotate(
            f"{kwh_fmt}\n{pct_fmt}",
            xy=(x * 0.68, y * 0.68),
            xytext=(x * 1.18, y * 1.18),
            ha=ha, va="center",
            fontsize=7.5, color="#333333",
            arrowprops=dict(arrowstyle="-", color="#888888", lw=0.8),
        )

    ax.set_aspect("equal")

    # Titre
    title = f"Répartition de la consommation\ndu site {client_name} en {saison_label}"
    ax.set_title(title, fontsize=8.5, color="#333333", pad=8)

    # Légende
    from matplotlib.patches import Patch
    legend_elements = [
        Patch(facecolor=_DONUT_COMPLEMENT, label="Complément réseau"),
        Patch(facecolor=_DONUT_ACC,        label="Autoconsommation Collective"),
    ]
    ax.legend(handles=legend_elements, loc="lower center",
              bbox_to_anchor=(0.5, -0.18), ncol=2, fontsize=7.5,
              frameon=False, handlelength=1.2)

    # Taux en bas
    fig.text(0.5, 0.01,
             f"Taux d'autoproduction {saison_label} : {taux:.0f} %",
             ha="center", fontsize=9, fontweight="bold", color="#333333")

    fig.tight_layout(rect=[0, 0.08, 1, 1])
    return _fig_bytes(fig)


# ── Mise à jour des tableaux ───────────────────────────────────────────────────

def _update_tariff_table(slide, tariff_option: str, tariff_rates: dict):
    """Met à jour le tableau des tarifs EDF (slide 7)."""
    for shape in slide.shapes:
        if shape.shape_type != 19:
            continue
        tbl = shape.table
        n_rows = len(tbl.rows)
        n_cols = len(tbl.columns)
        if n_rows < 2 or n_cols < 2:
            continue
        if tariff_option == "Monomial":
            v = tariff_rates.get("mono", 0)
            vals = [v, v, v, v]
        elif tariff_option == "HP / HC":
            hp = tariff_rates.get("hp", 0)
            hc = tariff_rates.get("hc", 0)
            vals = [hc, hc, hp, hp]   # HCE, HCH, HPE, HPH
        else:
            vals = [
                tariff_rates.get("hce", 0),
                tariff_rates.get("hch", 0),
                tariff_rates.get("hpe", 0),
                tariff_rates.get("hph", 0),
            ]
        for c, v in enumerate(vals[:n_cols - 1], 1):
            _set_cell_text(tbl.cell(1, c), f"{v:.1f}")
        break


def _update_economics_table(slide, eco: dict, client_name: str):
    """Met à jour le tableau des économies (slide 9)."""
    for shape in slide.shapes:
        if shape.shape_type != 19:
            continue
        tbl   = shape.table
        n_rows = len(tbl.rows)

        # En-tête : remplacer le nom du client
        if n_rows > 0 and len(tbl.columns) > 1:
            _replace_in_shape(shape, {"Roger Jalenques": client_name})

        # Lignes de données
        components = [
            ("Éco. Fourniture HT",    eco.get("dont Fourniture HT (€)", 0)),
            ("Éco. GC / Capacité HT", eco.get("dont GC HT (€)", 0)),
            ("Éco. CEE HT",           eco.get("dont CEE HT (€)", 0)),
            ("Éco. Accises HT",       eco.get("dont Accises HT (€)", 0) +
                                      eco.get("dont Autres taxes HT (€)", 0)),
            ("Économie totale HT",    eco.get("Économie grâce à l'ACC HT (€)", 0)),
        ]
        # On ne garde que les lignes avec une valeur non nulle
        components = [(k, v) for k, v in components if abs(v) > 0.01]

        for r, (label, value) in enumerate(components, 1):
            if r >= n_rows:
                break
            _set_cell_text(tbl.cell(r, 0), label)
            _set_cell_text(tbl.cell(r, 1), f"{value:,.0f} €")
        break


# ── Fonction principale ────────────────────────────────────────────────────────

def generate_pptx(
    template_path: str,
    client_name: str,
    location: str,
    project_type: str,
    date_debut: str,
    date_validite: str,
    annual: dict,
    monthly: pd.DataFrame,
    seasonal: pd.DataFrame,
    annual_eco: dict,
    consumer_keys: list,
    tariff_option: str,
    tariff_rates: dict,
    prix_cession: float,
    duree_contrat: int,
    p_start: pd.Timestamp,
    p_end: pd.Timestamp,
) -> bytes:
    """
    Génère et retourne les bytes du PPTX de proposition.
    """
    prs = Presentation(template_path)

    # ── Calculs préliminaires ──────────────────────────────────────────────────
    prod_gwh  = annual["Prod (kWh)"] / 1_000_000
    year_str  = (str(p_start.year)
                 if p_start.year == p_end.year
                 else f"{p_start.year}–{p_end.year}")
    period_str = f"{p_start.strftime('%d/%m/%Y')} → {p_end.strftime('%d/%m/%Y')}"

    taux_hiver = (seasonal.loc["Hiver (nov–mars)", "Taux ACC conso (%)"]
                  if "Hiver (nov–mars)" in seasonal.index else 0)
    taux_ete   = (seasonal.loc["Été (avr–oct)",    "Taux ACC conso (%)"]
                  if "Été (avr–oct)"    in seasonal.index else 0)

    k0  = consumer_keys[0] if consumer_keys else None
    eco = annual_eco.get(k0, {}) if k0 else {}

    duree_str = f"{duree_contrat} an{'s' if duree_contrat > 1 else ''}"

    # ── Remplacements textuels globaux ─────────────────────────────────────────
    common = {
        "Roger Jalenques":         client_name,
        "Saint-Julien-de-Toursac": location,
        "Centrale hydroélectrique": project_type,
        "1 janvier 2026":          date_debut,
        "30/09/2025":              date_validite,
        "5 ans":                   duree_str,
        " 140 ":                   f" {prix_cession:.0f} ",
        "140 €/MWh":               f"{prix_cession:.0f} €/MWh",
        "Te = 140":                f"Te = {prix_cession:.0f}",
    }
    for slide in prs.slides:
        _replace_all(slide, common)

    # ── Slide 3 : Production annuelle ──────────────────────────────────────────
    slide3 = prs.slides[2]
    _replace_all(slide3, {
        "2,3 GWh": f"{prod_gwh:.2f} GWh",
        "2024":     year_str,
    })
    img = chart_prod_mensuelle(monthly, w=6.66, h=4.12)
    _swap_picture_at(slide3, left_in=1.46, top_in=0.90, img_bytes=img)

    # ── Slide 5 : Hypothèses ──────────────────────────────────────────────────
    slide5 = prs.slides[4]
    _replace_all(slide5, {
        "1 janvier 2024 → 31 décembre 2024": period_str,
        "2024":                               year_str,
        "2,3 GWh":                            f"{prod_gwh:.2f} GWh",
    })

    # ── Slide 6 : Répartition de la consommation (donuts saisonniers) ────────
    slide6 = prs.slides[5]

    acc_hiver   = (seasonal.loc["Hiver (nov–mars)", "ACC (kWh)"]
                   if "Hiver (nov–mars)" in seasonal.index else 0)
    conso_hiver = (seasonal.loc["Hiver (nov–mars)", "Conso (kWh)"]
                   if "Hiver (nov–mars)" in seasonal.index else 0)
    acc_ete     = (seasonal.loc["Été (avr–oct)", "ACC (kWh)"]
                   if "Été (avr–oct)" in seasonal.index else 0)
    conso_ete   = (seasonal.loc["Été (avr–oct)", "Conso (kWh)"]
                   if "Été (avr–oct)" in seasonal.index else 0)

    img_hiver = chart_donut_saison(acc_hiver,  conso_hiver, client_name, "hiver", w=4.75, h=2.93)
    img_ete   = chart_donut_saison(acc_ete,    conso_ete,   client_name, "été",   w=4.75, h=2.93)
    _swap_picture_at(slide6, left_in=0.50, top_in=1.06, img_bytes=img_hiver)
    _swap_picture_at(slide6, left_in=5.13, top_in=1.06, img_bytes=img_ete)

    # ── Slide 7 : Proposition tarifaire ───────────────────────────────────────
    _update_tariff_table(prs.slides[6], tariff_option, tariff_rates)

    # ── Slide 9 : Calcul économies ────────────────────────────────────────────
    if eco:
        _update_economics_table(prs.slides[8], eco, client_name)

    # ── Export ─────────────────────────────────────────────────────────────────
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
